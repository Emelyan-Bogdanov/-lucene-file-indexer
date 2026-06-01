import csv
import json
import os
import re
import xml.etree.ElementTree as ET
from io import StringIO


def parse_txt(filepath):
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse_pdf(filepath):
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    except ImportError:
        return ""
    except Exception:
        return ""


def parse_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        return ""
    except Exception:
        return ""


def parse_xlsx(filepath):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                texts.append(" ".join(str(cell) for cell in row if cell is not None))
        return "\n".join(texts)
    except ImportError:
        return ""
    except Exception:
        return ""


def parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        return ""
    except Exception:
        return ""


def parse_html(filepath):
    try:
        from bs4 import BeautifulSoup
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n")
    except ImportError:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = re.sub(r"<[^>]+>", " ", f.read())
            return re.sub(r"\s+", " ", text).strip()
    except Exception:
        return ""


def parse_json(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return ""


def parse_xml(filepath):
    try:
        tree = ET.parse(filepath)
        root = tree.getroot()
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
        return "\n".join(texts)
    except Exception:
        return ""


def parse_source(filepath):
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse_csv(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            return "\n".join(" ".join(row) for row in reader)
    except Exception:
        return ""


PARSERS = {
    ".txt": parse_txt,
    ".md": parse_txt,
    ".rst": parse_txt,
    ".py": parse_source,
    ".js": parse_source,
    ".ts": parse_source,
    ".jsx": parse_source,
    ".tsx": parse_source,
    ".java": parse_source,
    ".c": parse_source,
    ".cpp": parse_source,
    ".h": parse_source,
    ".hpp": parse_source,
    ".cs": parse_source,
    ".rb": parse_source,
    ".go": parse_source,
    ".rs": parse_source,
    ".swift": parse_source,
    ".kt": parse_source,
    ".scala": parse_source,
    ".php": parse_source,
    ".pl": parse_source,
    ".sh": parse_source,
    ".bat": parse_source,
    ".ps1": parse_source,
    ".sql": parse_source,
    ".html": parse_html,
    ".htm": parse_html,
    ".json": parse_json,
    ".xml": parse_xml,
    ".yaml": parse_txt,
    ".yml": parse_txt,
    ".toml": parse_txt,
    ".ini": parse_txt,
    ".cfg": parse_txt,
    ".conf": parse_txt,
    ".csv": parse_csv,
    ".tsv": parse_csv,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
}

SOURCE_CODE_EXTS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".rb", ".go", ".rs", ".swift", ".kt", ".scala", ".php", ".pl",
    ".sh", ".bat", ".ps1", ".sql",
}


def get_parser(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    return PARSERS.get(ext, None)


def can_parse(filepath):
    return get_parser(filepath) is not None


def extract_text(filepath):
    parser = get_parser(filepath)
    if parser is None:
        return ""
    try:
        return parser(filepath)
    except Exception:
        return ""
